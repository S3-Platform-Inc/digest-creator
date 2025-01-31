import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from deep_translator import GoogleTranslator
from os.path import join
from tqdm import tqdm
from src.utils import current_time
import numpy as np
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData


class DigestCreator():
    src_name_beautify = {
        'w3c': 'W3C',
        'jcb': 'JCB',
        'visa': 'VISA',
        'emvco': 'EMVCo',
        'paymentsdive': 'PaymentsDive',
        'mit': 'MIT',
        'bis': 'BIS',
        'payments-journal': 'PaymentsJournal',
        'retailloyalty': 'RetailLoyalty',
        'eba': 'European Banking Association',
        'businesswire': 'Businesswire',
        'techcrunch': 'TechCrunch',
        'nfcw': 'NFCW',
        'thepaypers': 'ThePaypers',
        'americanexpress': 'American Express',
        'ieee': 'IEEE',
        'eupay': 'European Payments Council',
        'iso20022': 'ISO20022',
        'fido': 'FIDO',
        'finextra': 'Finextra',
        'kpmg': 'KPMG',
        'rfc': 'RFC',
        'ecb': 'European Central Bank',
        'pci': 'PCI',
        'eucommission': 'European Commission',
        'nist': 'NIST',
        'paypal': 'Paypal',
        'swift': 'SWIFT',
        'openbanking': 'OpenBanking',
        'openid': 'OpendID',
        'pwc': 'PWC'}

    def __init__(self):
        pass

    def strip_text(self, text: str, n_symbols: int = 500):
        stripped_text = ''

        for sent in text.split('.'):
            if len(stripped_text) <= n_symbols:
                stripped_text += sent
                stripped_text += '.'

        stripped_text = stripped_text[:-1]

        return stripped_text

    def add_title_slide(self, presentation: Presentation(), title: str = 'Дайджест инноваций'):
        slide_layout = presentation.slide_layouts[6]  # Using a blank slide layout
        slide = presentation.slides.add_slide(slide_layout)

        slide.shapes.add_picture("../SlideTemplates/top_bottom_blue_gradient_bk.png", 0, 0,
                                 width=presentation.slide_width,
                                 height=presentation.slide_height)

        images = [
            {'path': '../SlideTemplates/MIR_white_logo.png', 'left': Cm(1), 'top': Cm(1), 'height': Cm(0.6)},
            {'path': '../SlideTemplates/big_IT_leaf.png', 'left': Cm(16), 'top': Cm(3), 'height': Cm(21)},
            {'path': '../SlideTemplates/M_green_leaf.png', 'left': Cm(13.9), 'top': Cm(13.2), 'height': Cm(2.7),
             'rotation': 190},
            {'path': '../SlideTemplates/S_green_leaf.png', 'left': Cm(27.1), 'top': Cm(2.9), 'height': Cm(4),
             'rotation': 190},
            {'path': '../SlideTemplates/XS_green_leaf.png', 'left': Cm(22), 'top': Cm(2.9), 'height': Cm(1.4),
             'rotation': 190},

        ]

        for img in images:
            pic = slide.shapes.add_picture(img['path'], img['left'], img['top'], height=img['height'])
            if 'rotation' in img.keys():
                pic.rotation = img['rotation']

        textbox_title = slide.shapes.add_textbox(left=Cm(1), top=Cm(6.8), width=Cm(14), height=Cm(5.4))

        # Title with wrapping enabled
        textbox_title.text_frame.paragraphs[0].font.bold = True
        textbox_title.text_frame.paragraphs[0].font.size = Pt(60)
        textbox_title.text_frame.word_wrap = True
        textbox_title.text_frame.paragraphs[0].text = title

        textbox_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        textbox_info = slide.shapes.add_textbox(left=Cm(1), top=Cm(15.6), width=Cm(14), height=Cm(5.4))

        # Title with wrapping enabled
        textbox_info.text_frame.paragraphs[0].font.bold = False
        textbox_info.text_frame.paragraphs[0].font.size = Pt(20)
        textbox_info.text_frame.word_wrap = True
        textbox_info.text_frame.paragraphs[0].text = f'Департамент Инноваций, {current_time(return_str=False):%d.%m.%Y}'

        textbox_info.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    def add_graph_slide(self, presentation: Presentation(), df, title: str = 'Обзор источников'):

        result = df.groupby('fix_src_name')['to_digest'].value_counts().unstack(fill_value=0)

        # Convert to lists
        sources = result.index.tolist()
        src_0_count = result[0].tolist()  # Counts of 0
        src_1_count = result[1].tolist()  # Counts of 1

        slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # Using a blank slide layout

        # Define chart data
        chart_data = CategoryChartData()
        chart_data.categories = sources
        chart_data.add_series('Неинтересно', src_0_count)
        chart_data.add_series('Интересно', src_1_count)

        # Define position and size of the chart
        x, y, cx, cy = Cm(1), Cm(2.5), Cm(31.87), Cm(15.8)

        # Add the chart to the slide as a stacked bar chart
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_STACKED, x, y, cx, cy, chart_data
        )

        chart = graphic_frame.chart

        # Set colors for each series using points and add data labels
        for i, point in enumerate(chart.series[0].points):
            fill = point.format.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(2, 175, 255)
            point.has_data_labels = True  # Enable data labels
            data_label = point.data_label
            text_frame = data_label.text_frame
            # text_frame.clear()  # this line is not needed, assigning to .text does this
            text_frame.text = str.format(str(src_0_count[i]))
            # for paragraph in text_frame.paragraphs:
            #     paragraph.font.size = Pt(10)
            #
            # # -- OR --
            #
            # # for run in text_frame.paragraphs[0].runs:
            # #     run.font.size = Pt(10)
        for i, point in enumerate(chart.series[1].points):
            fill = point.format.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(31, 91, 215)
            point.has_data_labels = True  # Enable data labels
            point.data_label.text_frame.paragraphs[0].font.size = Pt(12)
            point.data_label.text_frame.paragraphs[0].text = str(src_1_count[i])  # Set label text
            # point.data_label.text_frame.paragraphs[0].font.size = Pt(12)
            # point.data_label.text_frame.paragraphs[0].font.color.rgb = RGBColor(90, 90, 90)

        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = Pt(14)

        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = Pt(14)
        value_axis.tick_labels.font.color.rgb = RGBColor(128, 128, 128)  # Grey color

        # Change grid lines color to grey (if applicable)
        value_axis.major_gridlines.format.line.color.rgb = RGBColor(128, 128, 128)  # Grey color

        # Change x-axis line color to grey
        value_axis.format.line.color.rgb = RGBColor(128, 128, 128)  # Grey color for x-axis line

        value_axis.format.line.width = Pt(0)

        # Change y-axis line color to grey (if applicable)
        category_axis.format.line.color.rgb = RGBColor(128, 128, 128)  # Grey color for y-axis line

        textbox_title = slide.shapes.add_textbox(left=Cm(1), top=Cm(1), width=Cm(14), height=Cm(1.2))

        # Title with wrapping enabled
        textbox_title.text_frame.paragraphs[0].font.bold = True
        textbox_title.text_frame.paragraphs[0].font.size = Pt(24)
        textbox_title.text_frame.word_wrap = True
        textbox_title.text_frame.paragraphs[0].text = title

        fill = textbox_title.text_frame.paragraphs[0].font.fill
        fill.gradient()  # Set fill type to gradient

        # Set up gradient stops
        gradient_stops = fill.gradient_stops

        # First stop (0% position)
        stop1 = gradient_stops[0]
        stop1.position = 0.0  # Position at start
        stop1.color.rgb = RGBColor(5, 167, 251)  # Red color

        # Second stop (100% position)
        stop2 = gradient_stops[1]
        stop2.position = 1.0  # Position at end
        stop2.color.rgb = RGBColor(29, 96, 217)  # Blue color

        # Optionally set the angle of the gradient
        fill.gradient_angle = 270  # Angle in degrees

    def add_doc_info_slide(self, presentation: Presentation(), df_row, translate_text: bool = False):
        # Create a new slide
        slide_layout = presentation.slide_layouts[6]  # Using a blank slide layout
        slide = presentation.slides.add_slide(slide_layout)

        # Add source name above the title
        source = df_row['fix_src_name']
        title = df_row['title']
        abstract = self.strip_text(text=df_row['abstract'], n_symbols=700)
        if source == 'W3C':
            try:
                abstract = abstract.split('Abstract')[1]
            except:
                pass

        web_link = df_row['weblink']

        user_comment = ''

        if df_row['user_1_comment'] == 'nan' and df_row['user_2_comment'] == 'nan' and df_row[
            'user_3_comment'] == 'nan':
            user_comment = 'Комментарий'
        elif not df_row['user_1_comment'] == 'nan':
            user_comment = user_comment + df_row['user_1_comment']
        elif not df_row['user_2_comment'] == 'nan':
            user_comment = user_comment + df_row['user_2_comment']
        elif not df_row['user_3_comment'] == 'nan':
            user_comment = user_comment + df_row['user_3_comment']
        else:
            raise Exception("Unsupported df_row['user_X_comment'] format!")

        # Define widths and heights for text boxes
        half_width = Cm(15)  # Half of the slide width
        abstract_width = Cm(15)
        title_width = Cm(25)
        title_height = Cm(2)  # Height for title box
        abstract_height = Cm(3)  # Height for abstract box

        # Add text boxes for source, title, abstract, and web link
        textbox_source = slide.shapes.add_textbox(Cm(1), Cm(1), half_width, Cm(1))

        fill = textbox_source.text_frame.paragraphs[0].font.fill
        fill.gradient()  # Set fill type to gradient

        # Set up gradient stops
        gradient_stops = fill.gradient_stops

        # First stop (0% position)
        stop1 = gradient_stops[0]
        stop1.position = 0.0  # Position at start
        stop1.color.rgb = RGBColor(217, 217, 217)  # Red color

        # Second stop (100% position)
        stop2 = gradient_stops[1]
        stop2.position = 1.0  # Position at end
        stop2.color.rgb = RGBColor(166, 166, 166)  # Blue color

        # Optionally set the angle of the gradient
        fill.gradient_angle = 270  # Angle in degrees

        # Set text for each text box
        textbox_source.text_frame.paragraphs[0].font.bold = True
        textbox_source.text_frame.paragraphs[0].font.size = Pt(20)
        textbox_source.text_frame.word_wrap = True
        textbox_source.text_frame.paragraphs[0].text = source

        textbox_title = slide.shapes.add_textbox(Cm(1), Cm(2), title_width, title_height)

        if len(title) <= 60:
            textbox_abstract = slide.shapes.add_textbox(Cm(1), Cm(4), abstract_width, abstract_height)
        elif len(title) <= 110:
            textbox_abstract = slide.shapes.add_textbox(Cm(1), Cm(5), abstract_width, abstract_height)
        else:
            textbox_abstract = slide.shapes.add_textbox(Cm(1), Cm(5), abstract_width, abstract_height)

        textbox_link = slide.shapes.add_textbox(Cm(1), Cm(16), Cm(5), Cm(2))
        slide_number = slide.shapes.add_textbox(Cm(32.5), Cm(17.8), Cm(1), Cm(0.7))
        rectangle_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Cm(18),
            Cm(15),
            Cm(14),
            Cm(6)
        )
        text_box_user_comment = slide.shapes.add_textbox(Cm(18.5), Cm(15.5), Cm(13), Cm(3))



        # Title with wrapping enabled
        textbox_title.text_frame.paragraphs[0].font.bold = True
        textbox_title.text_frame.paragraphs[0].font.size = Pt(28)
        textbox_title.text_frame.word_wrap = True
        textbox_title.text_frame.paragraphs[0].text = title

        # Abstract with wrapping enabled
        textbox_abstract.text_frame.paragraphs[0].text = abstract

        if len(abstract) <= 250:
            textbox_abstract.text_frame.paragraphs[0].font.size = Pt(16)
        else:
            textbox_abstract.text_frame.paragraphs[0].font.size = Pt(14)

        textbox_abstract.text_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)

        textbox_abstract.text_frame.add_paragraph()
        textbox_abstract.text_frame.add_paragraph()
        if translate_text:
            try:
                textbox_abstract.text_frame.paragraphs[2].text = GoogleTranslator(source='en', target='ru').translate(
                    abstract)
            except:
                textbox_abstract.text_frame.paragraphs[2].text = 'Перевод'
        else:
            textbox_abstract.text_frame.paragraphs[2].text = 'Перевод'
        textbox_abstract.text_frame.paragraphs[2].font.size = Pt(12)
        textbox_abstract.text_frame.paragraphs[2].font.color.rgb = RGBColor(191, 191, 191)

        textbox_abstract.text_frame.word_wrap = True

        # Add "More..." with hyperlink
        hyperlink_run = textbox_link.text_frame.paragraphs[0].add_run()
        hyperlink_run.text = 'Подробнее'
        hyperlink_run.hyperlink.address = web_link

        slide_number.text_frame.paragraphs[0].text = str(len(presentation.slides))
        slide_number.text_frame.paragraphs[0].font.size = Pt(12)
        slide_number.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        fill_format = rectangle_shape.fill
        fill_format.gradient()  # Apply gradient fill
        fill_format.gradient_angle = 135

        stops = fill_format.gradient_stops

        # Define colors for gradient stops (blue gradient)
        stops[0].color.rgb = RGBColor(31, 91, 215)
        stops[0].position = 0.0  # Start position of gradient stop

        stops[1].color.rgb = RGBColor(0, 175, 255)
        stops[1].position = 1.0  # End position of gradient stop

        # rectangle_shape.line.fill.transparency = 0.99

        text_box_user_comment.text_frame.paragraphs[0].text = user_comment
        text_box_user_comment.text_frame.word_wrap = True
        text_box_user_comment.text_frame.paragraphs[0].font.size = Pt(16)
        text_box_user_comment.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        pass

    def create_presentation(self, slide_format: str = 'wide'):
        presentation = Presentation()

        if slide_format == 'wide':
            presentation.slide_width = Cm(33.867)
            presentation.slide_height = Cm(19.05)
        else:
            raise Exception(f'{slide_format} is an unsupported format!')

        return presentation

    def save_presentation(self, presentation: Presentation(), name: str, path: str):
        presentation_file = join(path, f'{name}_{current_time()}.pptx')  # Output file name
        presentation.save(presentation_file)

        print(f"Презентация успешно создана: {presentation_file}")
        pass

    def load_data_table(self, df_path: str, df_name: str, sep: str = '\t'):

        df = pd.read_csv(filepath_or_buffer=join(df_path, df_name), sep=sep)

        df['src_name'] = df['src_name'].astype('str')
        df['abstract'] = df['abstract'].astype('str')
        df['user_1_score'] = df['user_1_score'].astype('Int64')
        df['user_2_score'] = df['user_2_score'].astype('Int64')
        df['user_3_score'] = df['user_3_score'].astype('Int64')
        df['user_1_comment'] = df['user_1_comment'].astype('str')
        df['user_2_comment'] = df['user_2_comment'].astype('str')
        df['user_3_comment'] = df['user_3_comment'].astype('str')

        df.drop_duplicates(subset=['title', 'published', 'weblink'], inplace=True)

        return df

    def create_digest(self, df_path, df_name, digest_name, save_path, df_sep: str = '\t', translate_text: bool = False,
                      score_sum_threshold: int = 1, exclude_from_min_1: list = None):

        df = self.load_data_table(df_path=df_path, df_name=df_name, sep=df_sep)

        df['sum_scores'] = df['user_1_score'].fillna(0) + df['user_2_score'].fillna(0) + df['user_3_score'].fillna(0)

        df['to_digest'] = np.where(df['sum_scores'] < score_sum_threshold, 0, 1)

        df['fix_src_name'] = df['src_name'].map(self.src_name_beautify)

        if exclude_from_min_1:
            print(f"Кол-во материалов, удовлетворяющих условию до изменения оценок: {df[df['to_digest'] == 1].shape[0]}")
            # Drop rows that have any of the values in the 'rebounds' column
            df.loc[df['weblink'].isin(exclude_from_min_1) & (df['to_digest'] == 1), 'to_digest'] = 0
            print(f'Изменено оценок с "1" на "0" после просмотра черновика: {len(exclude_from_min_1)}')

        df.sort_values(by='fix_src_name', inplace=True)

        if translate_text:
            print('Внимание! Включен перевод текста. Создание дайджеста займет больше времени.')
        presentation = self.create_presentation()

        self.add_title_slide(presentation=presentation)

        self.add_graph_slide(presentation=presentation, df=df, title='Обзор источников')

        df = df[df['to_digest'] == 1]

        print(f"Материал должен иметь от {score_sum_threshold} положительных оценок экспертов")
        print(f"Кол-во материалов, удовлетворяющих условию: {df.shape[0]}")

        for index, row in tqdm(df.iterrows(), total=df.shape[0]):
            try:
                self.add_doc_info_slide(presentation=presentation, df_row=row, translate_text=translate_text)
            except Exception as e:
                print(f'Возникла ошибке при обработке строки {row}')
                print(f'Ошибка:\n{e}')

        filename = f'{digest_name}_thr_{score_sum_threshold}_df_{df_name}'

        self.save_presentation(presentation=presentation, name=filename, path=save_path)
