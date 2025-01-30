from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.dml.color import RGBColor

from digest_creator import DigestCreator
import numpy as np

Digest = DigestCreator()
df = Digest.load_data_table(df_path="../Data/", df_name="view_experts_score_with_src_name_datetime.tsv", sep='\t')

df['sum_scores'] = df['user_1_score'].fillna(0) + df['user_2_score'].fillna(0) + df['user_3_score'].fillna(0)

# Assuming df is your DataFrame
df['min_1'] = np.where(df['sum_scores'] < 1, 0, 1)

print("df['min_1'].value_counts(dropna=False):")
print(df['min_1'].value_counts(dropna=False))

print("\ndf['src_name'].value_counts(dropna=False):")
print(df['src_name'].value_counts(dropna=False))

print("\ndf['src_name'].nunique()")
print(df['src_name'].nunique())

df.sort_values(by='src_name', inplace=True)

# Assuming df is your DataFrame
result = df.groupby('src_name')['min_1'].value_counts().unstack(fill_value=0)
print(result.info())
print(result)

# Convert to lists
sources = result.index.tolist()
src_0_count = result[0].tolist()  # Counts of 0
src_1_count = result[1].tolist()  # Counts of 1

# Display the lists
print("Sources:", sources)
print("Counts of 0:", src_0_count)
print("Counts of 1:", src_1_count)



# Sample data for the horizontal stacked bar chart with two categories
categories = sources
values_category_1 = src_0_count  # First category values
values_category_2 = src_1_count   # Second category values

# Create a presentation object
prs = Presentation()

# Add a slide to the presentation
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Using a blank slide layout

# Define chart data
chart_data = CategoryChartData()
chart_data.categories = categories
chart_data.add_series('Category 1', values_category_1)
chart_data.add_series('Category 2', values_category_2)

# Define position and size of the chart
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)

# Add the chart to the slide as a stacked bar chart
graphic_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_STACKED, x, y, cx, cy, chart_data
)

chart = graphic_frame.chart

# Set colors for each series using points and add data labels
for i, point in enumerate(chart.series[0].points):
    fill = point.format.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red for Category 1
    point.has_data_labels = True                 # Enable data labels
    point.data_label.text_frame.text = str(values_category_1[i])  # Set label text

for i, point in enumerate(chart.series[1].points):
    fill = point.format.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 255)  # Blue for Category 2
    point.has_data_labels = True                 # Enable data labels
    point.data_label.text_frame.text = str(values_category_2[i])  # Set label text

# Save the presentation
prs.save('../ChartTests/horizontal_stacked_bar_chart_with_data_labels.pptx')
