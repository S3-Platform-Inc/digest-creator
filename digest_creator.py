import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from deep_translator import GoogleTranslator
from os.path import join

from utils import current_time


class DigestCreator():

    def __init__(self):
        pass

    def create_and_format_slide(self, presentation: Presentation(), df_row, translate_text: bool = False):
        # Create a new slide
        slide_layout = presentation.slide_layouts[6]  # Using a blank slide layout
        slide = presentation.slides.add_slide(slide_layout)

        # Add source name above the title
        source = df_row['weblink'].split('/')[2]
        title = df_row['title']
        abstract = df_row['abstract']
        web_link = df_row['weblink']

        user_comment = ''

        if df_row['user_1_comment'] == 'nan' and df_row['user_2_comment'] == 'nan' and df_row[
            'user_3_comment'] == 'nan':
            user_comment = 'Комментарий'
        elif not df_row['user_1_comment'] == 'nan':
            user_comment = user_comment + df_row['user_1_comment']
        elif not df_row['user_2_comment'] == 'nan':
            user_comment = user_comment + df_row['user_2_comment']
        elif not df_row['user_3_comment'] == 'nan':
            user_comment = user_comment + df_row['user_3_comment']
        else:
            raise Exception("Unsupported df_row['user_X_comment'] format!")

        # Define widths and heights for text boxes
        half_width = Cm(15)  # Half of the slide width
        abstract_width = Cm(15)
        title_width = Cm(25)
        title_height = Cm(2)  # Height for title box
        abstract_height = Cm(3)  # Height for abstract box

        # Add text boxes for source, title, abstract, and web link
        textbox_source = slide.shapes.add_textbox(Cm(10), Cm(1), half_width, Cm(1))
        textbox_title = slide.shapes.add_textbox(Cm(1), Cm(2), title_width, title_height)
        textbox_abstract = slide.shapes.add_textbox(Cm(1), Cm(6), abstract_width, abstract_height)
        textbox_link = slide.shapes.add_textbox(Cm(1), Cm(16), Cm(5), Cm(2))
        slide_number = slide.shapes.add_textbox(Cm(32.5), Cm(17.8), Cm(1), Cm(0.7))
        rectangle_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Cm(18),
            Cm(15),
            Cm(14),
            Cm(6)
        )
        text_box_user_comment = slide.shapes.add_textbox(Cm(18.5), Cm(15.5), Cm(13), Cm(3))

        # Set text for each text box
        textbox_source.text = source

        # Title with wrapping enabled
        textbox_title.text_frame.paragraphs[0].font.bold = True
        textbox_title.text_frame.paragraphs[0].font.size = Pt(28)
        textbox_title.text_frame.word_wrap = True
        textbox_title.text_frame.paragraphs[0].text = title

        # Abstract with wrapping enabled
        textbox_abstract.text_frame.paragraphs[0].text = abstract
        textbox_abstract.text_frame.paragraphs[0].font.size = Pt(16)
        textbox_abstract.text_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)

        textbox_abstract.text_frame.add_paragraph()
        textbox_abstract.text_frame.add_paragraph()
        if translate_text:
            textbox_abstract.text_frame.paragraphs[2].text = GoogleTranslator(source='en', target='ru').translate(
                abstract)
        else:
            textbox_abstract.text_frame.paragraphs[2].text = 'Перевод'
        textbox_abstract.text_frame.paragraphs[2].font.size = Pt(12)
        textbox_abstract.text_frame.paragraphs[2].font.color.rgb = RGBColor(191, 191, 191)

        textbox_abstract.text_frame.word_wrap = True

        # Add "More..." with hyperlink
        hyperlink_run = textbox_link.text_frame.paragraphs[0].add_run()
        hyperlink_run.text = 'Подробнее'
        hyperlink_run.hyperlink.address = web_link

        slide_number.text_frame.paragraphs[0].text = str(len(presentation.slides))
        slide_number.text_frame.paragraphs[0].font.size = Pt(12)
        slide_number.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        fill_format = rectangle_shape.fill
        fill_format.gradient()  # Apply gradient fill
        fill_format.gradient_angle = 135

        stops = fill_format.gradient_stops

        # Define colors for gradient stops (blue gradient)
        stops[0].color.rgb = RGBColor(31, 91, 215)
        stops[0].position = 0.0  # Start position of gradient stop

        stops[1].color.rgb = RGBColor(0, 175, 255)
        stops[1].position = 1.0  # End position of gradient stop

        # rectangle_shape.line.fill.transparency = 0.99

        text_box_user_comment.text_frame.paragraphs[0].text = user_comment
        text_box_user_comment.text_frame.word_wrap = True
        text_box_user_comment.text_frame.paragraphs[0].font.size = Pt(16)
        text_box_user_comment.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        pass

    def create_presentation(self, slide_format: str = 'wide'):
        presentation = Presentation()

        if slide_format == 'wide':
            presentation.slide_width = Cm(33.867)
            presentation.slide_height = Cm(19.05)
        else:
            raise Exception(f'{slide_format} is an unsupported format!')

        return presentation

    def save_presentation(self, presentation: Presentation(), name: str, path: str):
        presentation_file = join(path, f'{name}_{current_time()}.pptx')  # Output file name
        presentation.save(presentation_file)

        print(f"Presentation created successfully: {presentation_file}")
        pass

    def load_data_table(self, df_path: str, df_name: str, sep: str = '\t'):

        df = pd.read_csv(filepath_or_buffer=join(df_path, df_name), sep=sep)

        df['sourceid'] = df['sourceid'].astype('str')
        df['abstract'] = df['abstract'].astype('str')
        df['user_1_score'] = df['user_1_score'].astype('Int64')
        df['user_2_score'] = df['user_2_score'].astype('Int64')
        df['user_3_score'] = df['user_3_score'].astype('Int64')
        df['user_1_comment'] = df['user_1_comment'].astype('str')
        df['user_2_comment'] = df['user_2_comment'].astype('str')
        df['user_3_comment'] = df['user_3_comment'].astype('str')

        return df

    def create_digest(self, df_path, df_name, digest_name, save_path, df_sep: str = '\t', translate_text: bool = False,
                      score_sum_threshold: int = 1):

        df = self.load_data_table(df_path=df_path, df_name=df_name, sep=df_sep)

        df['sum_scores'] = df['user_1_score'].fillna(0) + df['user_2_score'].fillna(0) + df['user_3_score'].fillna(0)

        df = df[df['sum_scores'] >= score_sum_threshold]

        presentation = self.create_presentation()

        for index, row in df.iterrows():
            self.create_and_format_slide(presentation=presentation, df_row=row, translate_text=translate_text)

        filename = f'{digest_name}_thr_{score_sum_threshold}_df_{df_name}'

        self.save_presentation(presentation=presentation, name=filename, path=save_path)
