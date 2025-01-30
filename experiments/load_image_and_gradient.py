from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum import shapes
# Create a presentation object
prs = Presentation()

# Add a blank slide
slide_layout = prs.slide_layouts[5]  # 5 is usually a blank slide
slide = prs.slides.add_slide(slide_layout)

# Define image paths and positions
images = [
    {'path': '../SlideTemplates/M_green_leaf.png', 'left': Inches(0.5), 'top': Inches(0.5), 'height': Inches(2)},
    {'path': '../SlideTemplates/big_IT_leaf.png', 'left': Inches(3), 'top': Inches(0.5), 'height': Inches(2)},
]

# Add images to the slide
for img in images:
    slide.shapes.add_picture(img['path'], img['left'], img['top'], height=img['height'])

# # Set background color (for example, light blue)
# background = slide.background
# fill = background.fill
# fill.solid()
# fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue color

slide.shapes.add_picture("../SlideTemplates/top_bottom_blue_gradient_bk.png", 0, 0, width=prs.slide_width, height=prs.slide_height)

# Save the presentation
prs.save('../../ImageTests/presentation_with_images_and_gradient.pptx')
