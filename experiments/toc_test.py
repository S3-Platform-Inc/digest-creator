from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

# Create a presentation object
prs = Presentation()

# Sample data for slides
slide_titles = ["Introduction", "Overview", "Details", "Conclusion"]

# Add slides with titles
for title in slide_titles:
    slide_layout = prs.slide_layouts[1]  # Using a title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

# Create a Table of Contents slide
toc_slide_layout = prs.slide_layouts[0]  # Using a title slide layout
toc_slide = prs.slides.add_slide(toc_slide_layout)
toc_title = toc_slide.shapes.title
toc_title.text = "Table of Contents"

# Add entries to the TOC with links to corresponding slides
for i, title in enumerate(slide_titles):
    left = Inches(1)
    top = Inches(1 + i * 0.5)  # Adjust top position for each entry
    width = Inches(8)
    height = Inches(0.5)

    # Create a text box for each TOC entry
    textbox = toc_slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame

    # Add the title as a paragraph and set it as a hyperlink
    p = text_frame.add_paragraph()
    p.text = title
    p.space_after = Pt(14)  # Space after paragraph

    # Create hyperlink to the corresponding slide
    r_id = toc_slide.part.relate_to(prs.slides[i + 1], RT.SLIDE)

    # Apply hyperlink to the text run
    r = p.add_run()
    r.text = title  # Set text for hyperlink
    r.hyperlink.address = None  # Clear any existing address if needed

    hyperlink_run = textbox_link.text_frame.paragraphs[0].add_run()
    hyperlink_run.text = 'Подробнее'
    hyperlink_run.hyperlink.address = web_link

# Save the presentation
prs.save('table_of_contents_with_links.pptx')
